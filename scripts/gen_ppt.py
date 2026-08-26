# -*- coding: utf-8 -*-
"""生成视觉化商策 PPT（B 赛道，16 页），含卡片/配色/图标/流程图 + 每页演讲备注"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 配色
NAVY = RGBColor(0x0E, 0x1A, 0x33)
NAVY_MID = RGBColor(0x1F, 0x35, 0x60)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
GOLD_LIGHT = RGBColor(0xE8, 0xC9, 0x6A)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
AMBER = RGBColor(0xC9, 0x9A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xB0, 0xBE, 0xD4)

# 卡片底色
CARD = RGBColor(0xF4, 0xF7, 0xFC)
CARD_GOLD = RGBColor(0xFB, 0xF6, 0xE8)
CARD_TEAL = RGBColor(0xE8, 0xF6, 0xF3)
CARD_CORAL = RGBColor(0xFD, 0xEF, 0xEB)
CARD_NAVY = RGBColor(0xE8, 0xED, 0xF7)

FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW, SH = Inches(13.333), Inches(7.5)


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, s, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = s.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        p.line_spacing = spacing
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.color.rgb = color
            r.font.bold = bold
    return tb


def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s


def header(slide, num, title):
    rect(slide, 0, 0, SW, Inches(1.0), NAVY)
    rect(slide, 0, Inches(1.0), SW, Pt(3), GOLD)
    text(slide, Inches(0.45), Inches(0.14), Inches(0.9), Inches(0.7), num, 28, GOLD, bold=True)
    text(slide, Inches(1.25), Inches(0.18), Inches(11.6), Inches(0.7), title, 23, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)


def card(slide, x, y, w, h, fill, accent, lines):
    """lines: list of (text, size, color, bold)"""
    shp = rect(slide, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE, line=accent)
    shp.adjustments[0] = 0.08
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (s, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = s
        p.space_after = Pt(5)
        p.line_spacing = 1.15
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.color.rgb = color
            r.font.bold = bold
    return shp


def badge(slide, x, y, d, s, bg, fg=WHITE, size=16):
    rect(slide, x, y, d, d, bg, MSO_SHAPE.OVAL)
    text(slide, x, y, d, d, s, size, fg, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ 1 封面 ============
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.7), SW, Pt(4), GOLD)
rect(s, Inches(0.8), Inches(5.0), Inches(3.2), Pt(2), NAVY_MID)
badge(s, Inches(1.0), Inches(1.0), Inches(1.0), "🛡", GOLD, NAVY, 34)
text(s, Inches(0.95), Inches(2.2), Inches(11.4), Inches(1.8), "合规智盾 Agent", 46, WHITE, bold=True)
text(s, Inches(0.95), Inches(3.9), Inches(11.4), Inches(0.7), "基于 AI Agent 的金融合规智能自检与申报系统", 21, GOLD_LIGHT, bold=True)
text(s, Inches(0.95), Inches(5.2), Inches(11.4), Inches(1.2), "暑假不回家队\n国元证券 · CFA 商业策划大赛 B 赛道", 15, LIGHT)
notes(s, "【开场 30 秒】语气自信平稳。报队名+题目，第一句点题：'我们把 CFA 道德准则，变成一线员工随开随用的合规 Agent。' 强调 'Agent' 这个主题关键词。")

# ============ 2 执行摘要 ============
s = prs.slides.add_slide(BLANK)
header(s, "01", "执行摘要：一句话 + 三个亮点")
text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.6), "员工用一句话描述行为 → Agent 自动识别准则、评估风险、执行申报", 20, NAVY, bold=True)
c1 = card(s, Inches(0.6), Inches(2.2), Inches(3.9), Inches(3.6), CARD, NAVY_MID, [
    ("🎯 不是问答工具", 17, NAVY, True),
    ("而是能「执行任务」的 Agent", 13, TEXT, False),
    ("", 10, TEXT, False),
    ("从「提建议」升维到「替办事」，补齐合规最缺的行动闭环", 12, MUTED, False),
])
c2 = card(s, Inches(4.7), Inches(2.2), Inches(3.9), Inches(3.6), CARD_TEAL, TEAL, [
    ("📖 可解释", 17, NAVY, True),
    ("RAG 约束，每个结论", 13, TEXT, False),
    ("都有题库案例作依据", 13, TEXT, False),
    ("", 10, TEXT, False),
    ("不输出「黑箱」判断", 12, MUTED, False),
])
c3 = card(s, Inches(8.8), Inches(2.2), Inches(3.9), Inches(3.6), CARD_GOLD, AMBER, [
    ("🔒 可落地", 17, NAVY, True),
    ("离线优先 + 邮件触达", 13, TEXT, False),
    ("", 10, TEXT, False),
    ("数据不出本机，国内 SMTP 真实发信，29 项测试全过", 12, MUTED, False),
])
notes(s, "【45 秒】总纲页，三个亮点先亮出来，后面逐页展开。核心一句话：'我们的核心，是把合规从提建议变成替办事。'")

# ============ 3 灵感来源 ============
s = prs.slides.add_slide(BLANK)
header(s, "02", "灵感来源：从一道题想到一个 Agent")
rect(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.6), CARD_GOLD, MSO_SHAPE.ROUNDED_RECTANGLE, line=GOLD)
text(s, Inches(0.9), Inches(1.65), Inches(11.6), Inches(1.3), "备考 CFA 道德时的真实观察：\n错题不在「记不住条文」，而在「具体场景如何对号入座」", 17, NAVY, bold=True)
text(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(0.5), "核心洞察", 16, GOLD, bold=True)
text(s, Inches(0.6), Inches(3.9), Inches(12.1), Inches(1.0), "道德合规的难点，从来不是「知识」，而是「判断」与「行动」", 20, NAVY, bold=True)
c = card(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.8), CARD, NAVY_MID, [
    ("结合大赛主题：AI 从「回答问题」跃迁为「执行任务」，恰好补上合规缺失的一环", 14, TEXT, False),
    ("于是：把 CFA 准则，变成一个能替你「做完该做的」的 Agent", 15, NAVY, True),
])
notes(s, "【50 秒】故事性最强，讲生动。'我们做 230 道题时发现，错的不是背准则，是不知道怎么套场景——这背后是个真问题。' 把「备考点」升华为「行业痛点」，体现灵感真实接地气。")

# ============ 4 切题程度 ============
s = prs.slides.add_slide(BLANK)
header(s, "03", "切题程度：三层严格对应")
c1 = card(s, Inches(0.6), Inches(1.5), Inches(3.9), Inches(5.2), CARD, NAVY_MID, [
    ("① 场景对应", 17, NAVY, True),
    ("三大场景锁定「合规」", 13, TEXT, False),
    ("", 10, TEXT, False),
    ("规则解析 / 制度比对 / 异常识别 / 报告生成 —— 正是系统四大功能", 12, MUTED, False),
])
c2 = card(s, Inches(4.7), Inches(1.5), Inches(3.9), Inches(5.2), CARD_TEAL, TEAL, [
    ("② 能力对应", 17, NAVY, True),
    ("Agent 四大核心能力", 13, TEXT, False),
    ("", 10, TEXT, False),
    ("自主感知 / 分析决策 / 任务执行 / 持续优化 —— 逐条落地", 12, MUTED, False),
])
c3 = card(s, Inches(8.8), Inches(1.5), Inches(3.9), Inches(5.2), CARD_GOLD, AMBER, [
    ("③ 标准对应", 17, NAVY, True),
    ("方案四大标准", 13, TEXT, False),
    ("", 10, TEXT, False),
    ("准确性 / 可解释性 / 安全性 / 落地能力 —— 逐条满足", 12, MUTED, False),
])
notes(s, "【50 秒】主动回应'切题吗'。'赛道给了场景、能力、标准三把尺子，我们用这三把尺子量自己，每一把都对得上。' 评委可能追问，提前背熟三层对应。")

# ============ 5 痛点 ============
s = prs.slides.add_slide(BLANK)
header(s, "04", "合规管理的三大痛点")
for i, (num, t, d) in enumerate([
    ("1", "判断门槛高", "准则表述开放，一线员工判断因人、因时而异"),
    ("2", "披露义务易遗漏", "不知道该不该披露、向谁披露、披露什么"),
    ("3", "申报与留痕割裂", "传统工具只提示风险，不执行申报"),
]):
    x = Inches(0.6 + i * 4.1)
    card(s, x, Inches(1.5), Inches(3.9), Inches(2.6), CARD_CORAL if i == 0 else CARD, CORAL if i == 0 else NAVY_MID, [
        (t, 17, NAVY, True),
        ("", 8, TEXT, False),
        (d, 13, TEXT, False),
    ])
    badge(s, x + Inches(0.15), Inches(1.7), Inches(0.5), num, CORAL if i == 0 else NAVY_MID, size=16)
rect(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(1.6), CARD_NAVY, MSO_SHAPE.ROUNDED_RECTANGLE, line=NAVY_MID)
text(s, Inches(0.9), Inches(4.6), Inches(11.6), Inches(1.2), "数据佐证：230 题库中准则应用指南占 64.3%，错题集中在客户责任与雇主责任的案例辨析\n结论：员工不是不想合规，而是缺少可执行的决策工具", 15, TEXT, bold=False)
notes(s, "【45 秒】讲具体，让评委有画面感。'员工不是不想合规，是面对抽象准则，不知道该不该、向谁、披露什么。' 点出：'所以我们要做的不仅是提示，更是帮他办成。'")

# ============ 6 方案概述 ============
s = prs.slides.add_slide(BLANK)
header(s, "05", "方案概述：一个能「执行」的合规 Agent")
# 流程图：输入 -> Agent -> 输出
rect(s, Inches(0.6), Inches(1.9), Inches(2.6), Inches(2.2), CARD, MSO_SHAPE.ROUNDED_RECTANGLE, line=NAVY_MID)
text(s, Inches(0.6), Inches(2.1), Inches(2.6), Inches(1.8), "👤 员工\n一句话\n描述行为", 15, NAVY, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(3.5), Inches(1.9), Inches(5.8), Inches(2.2), NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(3.5), Inches(2.0), Inches(5.8), Inches(2.0), "🛡 合规智盾 Agent\n识别准则 · 评估风险 · 执行申报", 16, WHITE, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(9.6), Inches(1.9), Inches(3.1), Inches(2.2), CARD_TEAL, MSO_SHAPE.ROUNDED_RECTANGLE, line=TEAL)
text(s, Inches(9.6), Inches(2.1), Inches(3.1), Inches(1.8), "✅ 输出\n风险+准则\n清单+申报单", 14, NAVY, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 箭头
for ax in [Inches(3.25), Inches(9.3)]:
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, Inches(2.7), Inches(0.45), Inches(0.45))
    ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background(); ar.shadow.inherit = False
text(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(0.5), "产品形态", 16, GOLD, bold=True)
c = card(s, Inches(0.6), Inches(5.05), Inches(12.1), Inches(1.7), CARD, NAVY_MID, [
    ("PWA 四视图（自检 / 准则速查 / 历史 / 合规部）+ 申报·自证·审批弹窗", 14, TEXT, False),
    ("支持语音输入 · 口语转正式 · 离线可用 · 手机/电脑多端适配", 14, TEXT, False),
])
notes(s, "【45 秒】介绍产品形态，点到为止，细节留给演示页。'它长这样——一个随开随用的网页应用，手机电脑都能用，断网也能用。'")

# ============ 7 四大能力 ============
s = prs.slides.add_slide(BLANK)
header(s, "06", "Agent 四大核心能力 · 逐一落地")
caps = [
    ("🔍 自主感知", "口语转正式 + 9类意图识别\n关键词+正则+语义三级", CARD, NAVY_MID),
    ("🧠 分析决策", "风险评级 + 22子条款\n+ RAG案例 + 行动强度分级", CARD_TEAL, TEAL),
    ("⚡ 任务执行", "申报+发邮件+审批+回执\n+自证(手写签名+电子印章)", CARD_GOLD, AMBER),
    ("🔄 持续优化", "增量知识库 + 缓存\n+ 留痕 + 合规小贴士", CARD_CORAL, CORAL),
]
for i, (t, d, fill, acc) in enumerate(caps):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.5 + row * 2.5)
    card(s, x, y, Inches(5.9), Inches(2.2), fill, acc, [
        (t, 17, NAVY, True),
        (d, 14, TEXT, False),
    ])
notes(s, "【60 秒】★核心页。'赛道要求 Agent 四大能力，我们一个不落——感知、决策、执行、优化，而且全部真实可运行，不是 PPT 概念。' 每个能力稍停，让评委看到对应关系。")

# ============ 8 Agent 作用 ============
s = prs.slides.add_slide(BLANK)
header(s, "07", "Agent 起到了什么作用")
roles = [
    ("👤 对员工", "从「凭感觉」到「有依据」\n口语转正式+申报+自证", CARD, NAVY_MID),
    ("🏛 对合规部门", "从「被动审查」到「签名审批」\n独立审查视角+回执留痕", CARD_TEAL, TEAL),
    ("🏢 对机构", "合规知识沉淀统一底座\n降本增效", CARD_GOLD, AMBER),
    ("🌐 对行业", "证明「规则驱动+案例约束\n+任务执行」范式可行", CARD_CORAL, CORAL),
]
for i, (t, d, fill, acc) in enumerate(roles):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.5 + row * 2.5)
    card(s, x, y, Inches(5.9), Inches(2.2), fill, acc, [
        (t, 16, NAVY, True),
        (d, 13, TEXT, False),
    ])
notes(s, "【45 秒】讲价值升维。'关键在任务执行——一般的 AI 只给建议，我们的 Agent 替你把申报发出去、把审批做完、把留痕存下来。员工和合规部是两套独立视角，这就是多智能体协同。'")

# ============ 9 技术架构 ============
s = prs.slides.add_slide(BLANK)
header(s, "08", "技术架构：RAG 约束，结论有据可依")
steps = ["输入", "意图分类", "向量检索 Top-5", "结构化输出", "任务执行"]
for i, st in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    rect(s, x, Inches(2.0), Inches(2.2), Inches(1.0), NAVY if i != 4 else GOLD, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, Inches(2.0), Inches(2.2), Inches(1.0), st, 13, WHITE if i != 4 else NAVY, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 4:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.5 + i * 2.5 + 2.25), Inches(2.25), Inches(0.25), Inches(0.45))
        ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background(); ar.shadow.inherit = False
text(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(0.5), "三大关键决策", 16, GOLD, bold=True)
decisions = [
    ("RAG 约束而非裸 LLM", "案例为依据，抑制幻觉"),
    ("规则引擎兜底", "没网、没 API 也能用"),
    ("离线优先", "数据不出本机，安全可落地"),
]
for i, (t, d) in enumerate(decisions):
    x = Inches(0.6 + i * 4.1)
    card(s, x, Inches(3.95), Inches(3.9), Inches(1.7), CARD, NAVY_MID, [
        (t, 15, NAVY, True),
        (d, 13, MUTED, False),
    ])
text(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.7), "技术栈：FastAPI + ChromaDB + RAG + PWA + SMTP（QQ 邮箱）", 13, MUTED)
notes(s, "【50 秒】不讲技术名词，讲'为什么这么选'。'我们不裸用大模型，而是用 RAG 让每个结论有案例兜底；再加规则引擎，没网也能跑。都是冲着安全、可解释、可落地去的。'")

# ============ 10 现场演示 ============
s = prs.slides.add_slide(BLANK)
header(s, "09", "现场演示：从识别到邮件触达的闭环")
demo = [
    ("输入", "客户送我去打高尔夫，还承担我出差的机票和酒店费用"),
    ("① 拦截", "高风险判定 + 红色风险拦截 + 准则 I(B)/IV(B)/VI(A)"),
    ("② 引用", "引用题库案例 M3-Q30 会议差旅招待"),
    ("③ 申报", "生成带唯一编号的申报单"),
    ("④ 触达", "一键发送邮件 → 合规部门邮箱实时收到"),
]
y = Inches(1.4)
for i, (t, d) in enumerate(demo):
    card(s, Inches(0.6), y, Inches(12.1), Inches(0.92), CARD if i != 4 else CARD_TEAL, NAVY_MID if i != 4 else TEAL, [
        (f"{t}　{'' if False else ''}", 14, NAVY, True),
    ])
    text(s, Inches(2.1), y + Inches(0.18), Inches(10.4), Inches(0.6), d, 13, TEXT)
    badge(s, Inches(0.75), y + Inches(0.18), Inches(0.55), str(i + 1), NAVY_MID, size=14)
    y += Inches(1.02)
notes(s, "【90 秒】★记忆点，全场最值得花时间。边操作边讲，放慢节奏。最后展示邮箱实时收到，强调'从识别到行动的完整闭环'。网络不稳立即切提前录好的视频。")

# ============ 11 数据与实证 ============
s = prs.slides.add_slide(BLANK)
header(s, "10", "数据基础与实证结果")
stats = [("230", "道题库"), ("22", "子条款全覆盖"), ("8", "类行为"), ("161", "合规案例"), ("29", "项测试全过")]
for i, (num, label) in enumerate(stats):
    x = Inches(0.55 + i * 2.5)
    rect(s, x, Inches(1.7), Inches(2.25), Inches(1.9), NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, Inches(1.85), Inches(2.25), Inches(0.9), num, 34, GOLD_LIGHT, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, Inches(2.85), Inches(2.25), Inches(0.5), label, 13, WHITE, align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(3.95), Inches(12.1), Inches(0.5), "实证验证", 16, GOLD, bold=True)
cases = [
    ("收礼招待 → I(B)/IV(B)/VI(A)，引用 M3-Q30", CARD),
    ("额外报酬 → IV(B)/VI(C)，引用 M3-Q12", CARD_TEAL),
    ("无关输入 → 正确返回 uncertain，不误报", CARD_GOLD),
]
for i, (c, fill) in enumerate(cases):
    x = Inches(0.6 + i * 4.1)
    card(s, x, Inches(4.5), Inches(3.9), Inches(1.5), fill, NAVY_MID, [
        (c, 13, TEXT, False),
    ])
notes(s, "【45 秒】用数据证明'做得出来、测得过'。'数据来源权威、处理规范，29 项测试全过，边界情况也不误报。'")

# ============ 12 五大挑战 ============
s = prs.slides.add_slide(BLANK)
header(s, "11", "五大挑战 · 逐一工程化应对")
challenges = [
    ("模型幻觉", "RAG 约束 + 温度 0.3 + 规则引擎兜底"),
    ("数据安全", "离线优先，数据不出本机 + 私有化部署"),
    ("可解释性", "每个结论附带准则条款与题库案例"),
    ("权限管理", "员工提交→合规签名审批→回执留痕，二次确认"),
    ("责任边界", "辅助决策定位，高风险强制人工确认"),
]
y = Inches(1.4)
for i, (ch, sol) in enumerate(challenges):
    card(s, Inches(0.6), y, Inches(4.2), Inches(0.92), CARD_CORAL, CORAL, [(f"⚔ {ch}", 15, NAVY, True)])
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.9), y + Inches(0.21), Inches(0.55), Inches(0.5))
    ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background(); ar.shadow.inherit = False
    card(s, Inches(5.6), y, Inches(7.1), Inches(0.92), CARD, NAVY_MID, [(sol, 14, TEXT, False)])
    y += Inches(1.05)
notes(s, "【60 秒】★关键页，评委必问。'赛道提了五大挑战，我们不是回避，是每一个都有工程实现。' 逐条念，补一句'这些不是纸面承诺，是已经在系统里落地的东西'。和问答准备对好。")

# ============ 13 改进过程 ============
s = prs.slides.add_slide(BLANK)
header(s, "12", "改进过程：从工具到 Agent 的六次迭代")
steps = [
    ("① 数据筑基", "230 题结构化"),
    ("② 感知决策", "9类分类+评级"),
    ("③ 可解释", "RAG 案例检索"),
    ("④ 关键跃迁", "任务执行→Agent"),
    ("⑤ 完善闭环", "拦截+自证+转写"),
    ("⑥ 多智能体", "合规审批+回执"),
]
for i, (t, d) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(1.5 + row * 2.3)
    card(s, x, y, Inches(3.9), Inches(2.0), CARD_NAVY if i != 3 else CARD_GOLD, NAVY_MID if i != 3 else AMBER, [
        (t, 16, NAVY, True),
        (d, 13, MUTED, False),
    ])
    badge(s, x + Inches(0.12), y + Inches(0.12), Inches(0.5), str(i + 1), NAVY_MID if i != 3 else AMBER, size=14)
text(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6), "第 ④ 步是质变（工具→Agent），第 ⑥ 步落地多智能体协同（审查 Agent）", 15, NAVY, bold=True)
notes(s, "【50 秒】呼应'持续优化'。'我们经历了六次迭代：第四步加上任务执行，从工具变成 Agent；第六步上线合规审批，落地了多智能体协同。持续优化不是口号，是我们实际走过的路。'")

# ============ 14 使用价值与可推广 ============
s = prs.slides.add_slide(BLANK)
header(s, "13", "使用价值与可推广性")
text(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.5), "使用价值", 16, GOLD, bold=True)
card(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(1.2), CARD_TEAL, TEAL, [
    ("员工降风险 · 合规部降成本 · 机构建文化 —— 从「事后追责」转向「事前自检」", 15, NAVY, True),
])
text(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(0.5), "可推广性", 16, GOLD, bold=True)
exps = [
    ("横向", "CFA 准则 → 软美元/GIPS/反洗钱"),
    ("纵向", "同一架构复用于投资、风控"),
    ("多智能体", "自检/申报/审查/审计协同"),
    ("跨行业", "医药/法律/证券强监管复用"),
]
for i, (t, d) in enumerate(exps):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(3.95 + row * 1.5)
    card(s, x, y, Inches(5.9), Inches(1.3), CARD, NAVY_MID, [
        (t + "：", 14, NAVY, True),
    ])
    text(s, x + Inches(1.7), y + Inches(0.18), Inches(4.0), Inches(0.9), d, 13, TEXT)
notes(s, "【45 秒】展示想象空间。'它不止解决 CFA 合规，这套「规则驱动+案例约束+任务执行」的方法论，可复制到更广合规域，甚至投资和风控场景。'")

# ============ 15 商业模式 ============
s = prs.slides.add_slide(BLANK)
header(s, "14", "商业模式与落地路径")
tiers = [
    ("个人版", "本地离线自检 · 免费", CARD),
    ("团队版", "内网部署/审计/批量 · 订阅", CARD_TEAL),
    ("定制版", "接入机构规则 · 项目制", CARD_GOLD),
]
for i, (t, d, fill) in enumerate(tiers):
    x = Inches(0.6 + i * 4.1)
    card(s, x, Inches(1.6), Inches(3.9), Inches(2.0), fill, NAVY_MID, [
        (t, 17, NAVY, True),
        (d, 13, MUTED, False),
    ])
text(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.5), "四阶段落地", 16, GOLD, bold=True)
steps = ["MVP", "机构试点", "扩域", "规模化"]
for i, st in enumerate(steps):
    x = Inches(0.6 + i * 3.1)
    rect(s, x, Inches(4.55), Inches(2.8), Inches(0.8), NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, Inches(4.55), Inches(2.8), Inches(0.8), st, 14, WHITE, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.85), Inches(4.75), Inches(0.25), Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background(); ar.shadow.inherit = False
text(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.6), "收入来源：企业订阅 + 私有化授权 + 内容授权 + 培训增值", 13, MUTED)
notes(s, "【35 秒】简略带过，时间留给演示和挑战应对。'商业模式分三层，从免费引流到企业订阅和私有化部署。'")

# ============ 16 结论+团队 ============
s = prs.slides.add_slide(BLANK)
header(s, "15", "创新点 · 团队 · 感谢")
text(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.5), "三大创新点", 16, GOLD, bold=True)
inns = [
    ("Agent 范式下沉", "从机构级下沉到个人合规行为辅助"),
    ("规则驱动可解释", "案例引用，结论有据可依"),
    ("离线优先可落地", "数据不出本机，真正可部署"),
]
for i, (t, d) in enumerate(inns):
    x = Inches(0.6 + i * 4.1)
    card(s, x, Inches(1.9), Inches(3.9), Inches(1.7), CARD, NAVY_MID, [
        (t, 15, NAVY, True),
        (d, 12, MUTED, False),
    ])
text(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.5), "团队分工", 16, GOLD, bold=True)
card(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(1.2), CARD_NAVY, NAVY_MID, [
    ("韩天翼（队长）：统筹+架构+核心研发 ｜ 彭德东：数据+调研 ｜ 禹夏航：测试+展示", 14, NAVY, True),
])
text(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.8), "感谢各位评委老师，欢迎提问", 22, NAVY, bold=True, align=PP_ALIGN.CENTER)
notes(s, "【40 秒】收尾三连：创新点→团队→感谢。'我们的创新在于，把 Agent 从机构级下沉到个人合规辅助，用规则驱动保证可解释，用离线架构保证可落地。以上是暑假不回家队的汇报，感谢评委，欢迎提问。' 说完原地等提问。")

out = r"D:\cfa-compliance-assistant\商策PPT_暑假不回家队.pptx"
prs.save(out)
print("已生成:", out)
